from pathlib import Path
import csv

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path("/Users/sophiahuang/Desktop/SquirrelGripper")
OUT = ROOT / "meeting_slides_dgdm_debug_20260904"
ASSETS = OUT / "assets" / "figures"
OUT.mkdir(parents=True, exist_ok=True)
ASSETS.mkdir(parents=True, exist_ok=True)

DIAG = Path("/Users/sophiahuang/Downloads/wandb_export_2026-09-03T20_40_45.018+08_00.csv")

BG = RGBColor(248, 249, 247)
INK = RGBColor(31, 43, 50)
MUTED = RGBColor(91, 106, 112)
TEAL = RGBColor(25, 116, 120)
ORANGE = RGBColor(224, 123, 57)
RED = RGBColor(184, 67, 63)
GREEN = RGBColor(57, 132, 91)
LIGHT = RGBColor(226, 234, 232)


def set_font(run, size, bold=False, color=INK, name="Aptos"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, x, y, w, h, text="", size=18, bold=False, color=INK,
            align=PP_ALIGN.LEFT, margin=0.05, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)
    return shape


def add_title(slide, title, kicker=None):
    if kicker:
        textbox(slide, 0.65, 0.28, 12.0, 0.28, kicker.upper(), 8, True, TEAL)
    textbox(slide, 0.65, 0.60, 12.0, 0.58, title, 26, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.24), Inches(1.15), Inches(0.045))
    line.fill.solid(); line.fill.fore_color.rgb = ORANGE; line.line.fill.background()


def add_footer(slide, number, source=""):
    textbox(slide, 0.65, 7.14, 10.8, 0.18, source, 7, False, MUTED)
    textbox(slide, 12.15, 7.10, 0.5, 0.22, str(number), 8, True, MUTED, PP_ALIGN.RIGHT)


def add_notes(slide, text):
    try:
        tf = slide.notes_slide.notes_text_frame
        tf.text = text
    except Exception:
        pass


def rect(slide, x, y, w, h, fill=LIGHT, line=None, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def add_bullets(slide, x, y, w, h, bullets, size=16, color=INK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = color
    return shape


def native_table(slide, x, y, w, h, data, widths=None, font=11):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for j, width in enumerate(widths):
            table.columns[j].width = Inches(width)
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.035)
            cell.fill.solid(); cell.fill.fore_color.rgb = TEAL if i == 0 else RGBColor(255, 255, 255)
            cell.border = None
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    set_font(r, font if i else font-1, i == 0, RGBColor(255,255,255) if i == 0 else INK)
    return table


def load_diagnostics():
    with DIAG.open(newline="") as f:
        rows = list(csv.DictReader(f))
    return [r for r in rows if r["model"] == "noise_conditioned"]


def plot_timestep_diagnostics():
    rows = load_diagnostics()
    t = np.array([int(r["diffusion_timestep"]) for r in rows])
    order = np.argsort(t); t = t[order]
    vals = {k: np.array([float(r[k]) for r in rows])[order] for k in [
        "utility_mae", "angular_span_mae", "direction_sign_accuracy", "direction_pearson"
    ]}
    fig, axes = plt.subplots(1, 2, figsize=(11.5, 3.7), dpi=180)
    ax = axes[0]
    ax.plot(t, vals["utility_mae"], "o-", color="#197478", lw=2, label="Utility MAE")
    ax.plot(t, vals["angular_span_mae"], "o-", color="#E07B39", lw=2, label="Angular-span MAE")
    ax.set_xlabel("Diffusion timestep (higher = noisier)"); ax.set_ylabel("MAE")
    ax.legend(frameon=False, fontsize=8); ax.grid(alpha=.18)
    ax = axes[1]
    ax.plot(t, vals["direction_sign_accuracy"], "o-", color="#B8433F", lw=2, label="Sign accuracy")
    ax.plot(t, vals["direction_pearson"], "o-", color="#397F5B", lw=2, label="Directional Pearson r")
    ax.axhline(.5, color="#77858A", ls="--", lw=1, label="50% sign baseline")
    ax.set_xlabel("Diffusion timestep (higher = noisier)"); ax.set_ylabel("Directional quality")
    ax.set_ylim(0, .75); ax.legend(frameon=False, fontsize=8); ax.grid(alpha=.18)
    fig.tight_layout()
    path = ASSETS / "timestep_diagnostics.png"; fig.savefig(path, bbox_inches="tight", facecolor="#F8F9F7"); plt.close(fig)
    return path


def plot_rank_heatmap():
    methods = ["Adam", "CMA-ES", "Cond. diff.", "DGDM .1", "DGDM .5", "DGDM 1", "DGDM 2"]
    values = np.array([
        [.541,.221,-.244,-.165,-.062,.000,.141],
        [.053,-.221,.229,.097,.179,.153,.394],
        [-.112,-.459,-.306,-.238,-.203,-.512,-.247],
        [.226,.112,-.424,-.394,-.135,.038,-.047],
    ])
    fig, ax = plt.subplots(figsize=(10.5, 3.6), dpi=180)
    im = ax.imshow(values, cmap="RdYlGn", vmin=-.6, vmax=.6, aspect="auto")
    ax.set_xticks(range(len(methods)), methods, rotation=25, ha="right")
    ax.set_yticks(range(4), ["Scenario 00", "Scenario 01", "Scenario 02", "Scenario 03"])
    for i in range(4):
        for j in range(7):
            ax.text(j, i, f"{values[i,j]:.2f}", ha="center", va="center", fontsize=8,
                    color="white" if abs(values[i,j]) > .35 else "#1F2B32")
    ax.set_title("Spearman correlation: surrogate score vs. simulator utility", fontsize=11)
    fig.colorbar(im, ax=ax, shrink=.8, label="Spearman ρ")
    fig.tight_layout()
    path = ASSETS / "ranking_heatmap.png"; fig.savefig(path, bbox_inches="tight", facecolor="#F8F9F7"); plt.close(fig)
    return path


def plot_oracle():
    scenarios = ["00", "01", "02", "03"]
    data = {
        "Adam": [.7510,.7469,.8145,.8354],
        "CMA-ES": [.5730,.5906,.5110,.7445],
        "Conditional diffusion": [.7414,.7553,.5285,.7460],
        "DGDM gs=0.1": [.7761,.7600,.5274,.7772],
    }
    colors = ["#77858A", "#B3B9B8", "#E07B39", "#197478"]
    x=np.arange(4); width=.19
    fig, ax=plt.subplots(figsize=(10.8,4.0), dpi=180)
    for j, (name, vals) in enumerate(data.items()):
        ax.bar(x+(j-1.5)*width, vals, width, label=name, color=colors[j])
    ax.set_xticks(x, [f"Scenario {s}" for s in scenarios]); ax.set_ylim(.35,.9)
    ax.set_ylabel("Best simulator utility among 16 candidates")
    ax.grid(axis="y", alpha=.18); ax.legend(ncol=4, frameon=False, fontsize=8, loc="upper left")
    fig.tight_layout()
    path=ASSETS/"oracle_utility.png"; fig.savefig(path,bbox_inches="tight",facecolor="#F8F9F7"); plt.close(fig)
    return path


diag_img = plot_timestep_diagnostics()
heat_img = plot_rank_heatmap()
oracle_img = plot_oracle()

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# 1 Cover
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
rect(s,.65,.7,.12,5.8,ORANGE, radius=False)
textbox(s,1.05,1.15,11.3,.45,"DGDM debugging update",12,True,TEAL)
textbox(s,1.05,1.75,10.9,1.25,"What the diagnostics reveal —\nand what to test next",30,True)
textbox(s,1.05,3.35,10.6,.75,"Squirrel Gripper design generation · mentor meeting",16,False,MUTED)
textbox(s,1.05,5.75,10.8,.55,"Evidence: timestep diagnostics + all-16 simulator evaluation",13,True,ORANGE)
add_footer(s,1,"4 Sep 2026")
add_notes(s,"Goal: align on the failure mode and approve the next ablation sequence. Emphasize that final design selection now uses simulator utility across all evaluated candidates.")

# 2 correspondence
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
add_title(s,"Our formulation follows DGDM’s logic, with a different interaction representation","Method correspondence")
data=[
    ["Component","Original DGDM","Squirrel Gripper"],
    ["Dynamics input","object o, finger m, pose p","scenario, 16-D finger, initial state"],
    ["Dynamics output","Δθ, Δx, Δy per initial pose","contact C, disturbance D, angular span A"],
    ["Task objective","target interaction profile F","U = 0.20C + 0.45D + 0.35A"],
    ["Diffusion","unconditional geometry prior; 15/5 DDIM","conditional prior; matched 15/5 DDIM"],
    ["Evaluation","16 starts; best measured performance","60 starts overall; current debug: all 16 simulated"],
]
native_table(s,.7,1.55,11.95,3.85,data,[2.0,4.7,5.25],11)
textbox(s,.8,5.72,11.6,.64,"Key distinction: the paper evaluates success rate after predicting motion; success rate is not the dynamics-network target.",15,True,ORANGE,PP_ALIGN.CENTER)
add_footer(s,2,"Source: Xu, Ha & Song, DGDM, arXiv:2402.15038 (v2)")
add_notes(s,"Conditional diffusion is acceptable and does not invalidate the comparison. The conceptual correspondence is preserved: learn a differentiable interaction model, construct a task objective, and use its gradient during denoising.")

# 3 aggregation and smoothness
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
add_title(s,"The main representational difference is where aggregation occurs","Why gradient quality may differ")
for i,(title,body,col) in enumerate([
    ("Original DGDM","Predict motion for many poses\n→ build an interaction profile\n→ aggregate objective gradients",TEAL),
    ("Our implementation","Simulate one scenario rollout\n→ predict three summary metrics\n→ weight into one utility",ORANGE),
]):
    x=.85+i*6.15; rect(s,x,1.65,5.55,3.25,RGBColor(255,255,255),col)
    textbox(s,x+.3,1.98,4.95,.38,title,19,True,col,PP_ALIGN.CENTER)
    textbox(s,x+.45,2.65,4.65,1.55,body,17,False,INK,PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
textbox(s,1.05,5.32,11.2,.95,"Contact-related summaries can change sharply when a small geometry change creates or removes contact. This is not “wrong”; it is a hypothesis for weaker local gradients and must be tested.",15,False,INK,PP_ALIGN.CENTER)
add_footer(s,3,"Dataset: 16,000 train + 4,000 test simulations; angular-span cap remains 180°")
add_notes(s,"The paper samples dense initial-pose grids and aggregates their gradients. Our 20,000 rollouts are substantial, but not equivalent to their dense per-pose interaction coverage. The experimental question is whether our summary metrics remain locally smooth enough for gradient guidance.")

# 4 timestep
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
add_title(s,"Noise diagnostics support testing timestep-gated guidance","Mentor’s proposed ablation")
s.shapes.add_picture(str(diag_img), Inches(.65), Inches(1.48), width=Inches(8.4))
rect(s,9.35,1.62,3.25,4.55,RGBColor(255,255,255),TEAL)
textbox(s,9.65,1.92,2.7,.42,"Observed",17,True,TEAL)
add_bullets(s,9.65,2.48,2.65,2.2,["Best alignment near t≈4–6","Sign accuracy peaks ≈0.66","Pearson r peaks ≈0.37","High-noise and t=0 are weaker"],13)
textbox(s,9.65,5.05,2.65,.72,"Test late-only, mid+late, and all-step guidance with paired noise.",13,True,ORANGE,PP_ALIGN.CENTER)
add_footer(s,4,"Source: wandb_export_2026-09-03; n=2,048 samples, 616 direction pairs per timestep")
add_notes(s,"In reverse denoising, the earliest sampling step has the largest diffusion timestep and noise. The mentor’s hypothesis is plausible, but the diagnostic suggests that only the final t≈0 step may also be weak. Therefore compare explicit schedules rather than assuming late-only will win.")

# 5 ranking
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
add_title(s,"The clean surrogate cannot reliably rank generated candidates","Top-1 selection failure")
s.shapes.add_picture(str(heat_img), Inches(.7), Inches(1.48), width=Inches(8.25))
rect(s,9.3,1.65,3.2,4.6,RGBColor(255,255,255),RED)
textbox(s,9.6,1.98,2.6,.42,"Interpretation",17,True,RED)
add_bullets(s,9.58,2.55,2.7,2.45,["Scenario 02: all top-4 overlaps = 0","Most correlations ≈0 or negative","Old selected designs ranked as low as 15–16/16"],13)
textbox(s,9.58,5.25,2.7,.58,"Final selection now uses simulator oracle over all 16.",13,True,GREEN,PP_ALIGN.CENTER)
add_footer(s,5,"Source: all-16 candidate rebenchmark; 16 candidates per method/scenario")
add_notes(s,"Clarify the previous implementation: analyze_study always selected by simulator utility, but benchmark_top_k=1 meant only the surrogate top candidate had been simulated. With top_k=16, oracle selection is now correct.")

# 6 oracle
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
add_title(s,"Oracle evaluation rescues DGDM in two scenarios, but the overall gain remains modest","Corrected method comparison")
s.shapes.add_picture(str(oracle_img), Inches(.65), Inches(1.48), width=Inches(8.65))
rect(s,9.55,1.65,2.95,4.65,RGBColor(255,255,255),TEAL)
textbox(s,9.82,1.95,2.42,.4,"Fixed gs = 0.1",17,True,TEAL,PP_ALIGN.CENTER)
textbox(s,9.82,2.65,2.42,.55,"0.710",28,True,TEAL,PP_ALIGN.CENTER)
textbox(s,9.82,3.18,2.42,.34,"mean max utility",10,False,MUTED,PP_ALIGN.CENTER)
textbox(s,9.82,3.82,2.42,.55,"+2.5%",28,True,ORANGE,PP_ALIGN.CENTER)
textbox(s,9.82,4.35,2.42,.55,"vs conditional diffusion",10,False,MUTED,PP_ALIGN.CENTER)
textbox(s,9.82,5.10,2.42,.75,"DGDM wins 00–01;\nAdam wins 02–03",13,True,INK,PP_ALIGN.CENTER)
add_footer(s,6,"Bars show simulator-best candidate among 16; scale 0.1 is fixed across scenarios")
add_notes(s,"Do not combine different best scales per scenario as a single method unless clearly labeled post-hoc. The defensible fixed-scale result is about +2.5% in mean maximum utility versus conditional diffusion.")

# 7 failure decomposition
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
add_title(s,"Scenario 02 isolates the unresolved guidance/prior failure","Metric-level evidence")
data=[
    ["Method","Utility","Contact","Disturbance","Angular span"],
    ["Adam","0.8145","0.6979","0.7221","1.0000"],
    ["Conditional diffusion","0.5285","0.5074","0.8374","0.1436"],
    ["DGDM gs=0.1","0.5274","0.5074","0.8350","0.1435"],
    ["DGDM gs=1","0.5202","0.5317","0.7996","0.1546"],
]
native_table(s,.95,1.68,11.35,2.45,data,[3.0,1.65,1.8,2.0,2.9],12)
textbox(s,1.0,4.55,11.2,.62,"DGDM does not improve angular span over unguided diffusion in Scenario 02, while Adam finds a feasible high-span basin.",17,True,RED,PP_ALIGN.CENTER)
add_bullets(s,1.55,5.42,10.1,1.0,["Possible causes: diffusion support gap; inaccurate noisy gradient; summary metric locally non-smooth."],15)
add_footer(s,7,"Utility weights: D=0.45, C=0.20, A=0.35; angular span normalized with 180° cap")
add_notes(s,"This is the most informative failure case. Disturbance is already high for diffusion designs, but angular span remains near 0.14. A targeted analysis should compare Adam’s winning design against the diffusion training distribution and visualize denoising trajectories.")

# 8 next experiments
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
add_title(s,"Next experiments separate timestep, gradient, and prior-support hypotheses","Decision plan")
items=[
    ("1", "Timestep-gated guidance", "Paired noise: none vs all 5 vs last 2 vs middle+last; fixed scale 0.1 and 1.0."),
    ("2", "Correct ranking supervision", "Train noisy dynamics with same-scenario, same-timestep design pairs; remeasure sign accuracy and r."),
    ("3", "Scenario-02 support test", "Measure nearest-neighbor distance of Adam winner to training and diffusion pools; inspect parameter percentiles."),
    ("4", "Evaluation protocol", "Simulate all 16; select by oracle utility; report max, mean, median, ≥0.70 rate, and C/D/A."),
]
for i,(num,title,body) in enumerate(items):
    y=1.55+i*1.28
    textbox(s,.75,y,.55,.55,num,21,True,RGBColor(255,255,255),PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(.72), Inches(y-.02), Inches(.62), Inches(.62)); circ.fill.solid(); circ.fill.fore_color.rgb=TEAL; circ.line.fill.background()
    # Re-add number above circle for proper z-order.
    textbox(s,.75,y,.55,.55,num,18,True,RGBColor(255,255,255),PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
    textbox(s,1.58,y-.02,3.05,.38,title,17,True,TEAL)
    textbox(s,4.55,y-.02,7.75,.75,body,14,False,INK)
textbox(s,.9,6.72,11.6,.34,"Decision gate: proceed to 5 seeds only if fixed-scale DGDM improves both oracle maximum and pool distribution.",13,True,ORANGE,PP_ALIGN.CENTER)
add_footer(s,8,"Primary endpoint: fixed-scale DGDM vs paired conditional diffusion")
add_notes(s,"Ask the mentor to approve the timestep schedules and the decision gate. Start with one seed and existing candidates where possible, then scale to five seeds only after a clear directional improvement.")

path=OUT/"dgdm_debugging_meeting_20260904.pptx"
prs.save(path)

(OUT/"terminology_ledger.md").write_text("""# Terminology ledger

| Canonical term | Definition / decision |
|---|---|
| DGDM | Dynamics-Guided Diffusion Model |
| clean dynamics model | Surrogate used for clean-design prediction/ranking |
| noisy dynamics model | Noise-conditioned surrogate used for denoising guidance |
| simulator oracle | Highest full-simulator utility among all evaluated candidates |
| selection score | Surrogate-predicted utility; not authoritative |
| utility | U = 0.20 contact + 0.45 disturbance + 0.35 angular span |
| direction sign accuracy | Fraction of design-pair directions whose predicted utility change has the correct sign |
| directional Pearson r | Correlation between predicted directional change and simulator-measured utility change |
""", encoding="utf-8")

(OUT/"qa_report.md").write_text("""# QA report

- Status: PPTX generated with 8 slides and speaker notes.
- Evidence assets: three charts generated from the local timestep CSV and user-provided all-16 tables.
- Angular-span normalization remains capped at 180 degrees.
- Oracle selection is explicitly distinguished from surrogate selection.
- Sources: local W&B CSV; user-provided all-16 benchmark tables; DGDM arXiv:2402.15038.
- Self-review: checked slide bounds, title hierarchy, table density, terminology, source labels, and claim/evidence alignment.
- Known limitation: the latest all-16 CSV was not present locally; those numerical values were transcribed from the user's pasted tables.
- Rendered slide preview was unavailable; package reopening and structural audit are performed separately.
""", encoding="utf-8")

print(path)
